from __future__ import print_function
import pickle
import os.path
import tempfile
import logging
import os
import io
import time
from googleapiclient.discovery import build
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
from pptx import Presentation
from pptx.util import Inches
import mimetypes

# If modifying these scopes, delete the file token.pickle.
SCOPES = ['https://www.googleapis.com/auth/drive', 'https://www.googleapis.com/auth/admin.directory.user']

MIME_TYPES = {
    "ppt": "application/vnd.ms-powerpoint",
    "pptx": "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation",
    "doc": "application/msword",
    "docx_files": "application/vnd.openxmlformats-officedocument."
                  "wordprocessingml.document",
    "xls": "application/vnd.ms-excel",
    "xlsx": "application/vnd.openxmlformats-officedocument."
            "spreadsheetml.sheet",
    "pdf": "application/pdf",
    "gslide": "application/vnd.google-apps.presentation",
    "gdoc": "application/vnd.google-apps.document",
    "gsheet": "application/vnd.google-apps.spreadsheet",
}


class Document_XML:
    def __init__(self, mimetype):
        self.mime_type = mimetype


DOC_TYPES = {
    'docx_files': Document_XML,

}


def get_file(file_type):
    return DOC_TYPES[file_type]


def modify_file_with_path(service, file_id, path, logger):
    file_downloaded = download_file(service, file_id, path, logger)
    logger.debug("the file has been downloaded...")
    file_name = file_downloaded.split('/')[1]
    logger.debug("opening file a and editing...")
    file_mimetype = _get_local_file_mime_type(file_name)
    # MIME_TYPES
    logger.debug(file_mimetype)
    andre
    open_presentation_and_edit(file_downloaded)
    # mime_pptx = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    media = MediaFileUpload(file_downloaded,
                            mimetype=file_mimetype)

    result = service.files().update(fileId=file_id, media_body=media, fields="id").execute()
    logger.debug("gdrive file updated...")
    return result["id"]


def download_file(service, file_id, path, logger):
    if not os.path.exists(path):
        os.makedirs(path)

    request = service.files().get_media(fileId=file_id)
    name = service.files().get(fileId=file_id).execute()['name']
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while done is False:
        (status, done) = downloader.next_chunk()
        logger.debug("Download {}%".format(int(status.progress() * 100)))

    file_name = os.path.join(path, name)
    with open(file_name, 'wb') as f:
        f.write(fh.getvalue())

    return file_name


def open_presentation_and_edit(pptx):
    prs = Presentation(pptx)
    blank_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "This is text inside a textbox"

    prs.save(pptx)


def get_user_files():
    # Call the Drive v3 API
    results = get_service().files().list(
        pageSize=10, fields="nextPageToken, files(id, name)").execute()
    items = results.get('files', [])

    if not items:
        print('No files found.')
    # else:
    #     print('Files:')
    #     for item in items:
    #         print(u'{0} ({1})'.format(item['name'], item['id']))
    return items


def _get_local_file_mime_type(name):
    return mimetypes.guess_type(name)[0]


def main():
    """Shows basic usage of the Drive v3 API.
    Prints the names and ids of the first 10 files the user has access to.
    """
    creds = None
    # The file token.pickle stores the user's access and refresh tokens, and is
    # created automatically when the authorization flow completes for the first
    # time.
    if os.path.exists('token.pickle'):
        with open('token.pickle', 'rb') as token:
            creds = pickle.load(token)
    # If there are no (valid) credentials available, let the user log in.
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                'credentials.json', SCOPES)
            creds = flow.run_local_server(port=0)
        # Save the credentials for the next run
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)

    # create logger
    logger = logging.getLogger('simple_example')
    logger.setLevel(logging.DEBUG)

    # create console handler and set level to debug
    ch = logging.StreamHandler()
    ch.setLevel(logging.DEBUG)

    # create formatter
    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')

    # add formatter to ch
    ch.setFormatter(formatter)

    # add ch to logger
    logger.addHandler(ch)
    service = get_service(creds)
    get_user_files(service)
    file_gdoc = '15yPu8XthjQeETTxj1STxZAkCoRqtKnkrMnW0jvWdKv8'
    # file_id = "1yCay0KFCjaS-k1fLROqsI9KWUUocl4KR"
    logger.debug(modify_file_with_path(service, file_gdoc, "generated", logger))


def get_service():
    from project.google_drive import obtain_creds
    return build('drive', 'v3', credentials=obtain_creds())


def get_user_file(file_id):
    result = get_service().files().get(
        fileId=file_id, fields="id, name").execute()

    if not result:
        print('No file found.')
    return result


if __name__ == '__main__':
    main()
